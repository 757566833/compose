import subprocess
import os
import io
import hashlib
import zipfile
import mimetypes
import tempfile
import shutil
import magic
import glob
from typing import List, Dict, Any

from dotenv import load_dotenv # 引入加载函数

import boto3
from fastapi import FastAPI, HTTPException
from pptx import Presentation
from pydantic import BaseModel
from botocore.client import Config

app = FastAPI(title="PPT-Parse Service - High Performance Version")

load_dotenv()
# --- 从环境变量读取配置 ---
OSS_URL = os.getenv("OSS_URL", "http://localhost:9000")
OSS_ACCESS_KEY = os.getenv("OSS_ACCESS_KEY", "OSS_ACCESS_KEY")
OSS_SECRET_KEY = os.getenv("OSS_SECRET_KEY", "OSS_SECRET_KEY")
OSS_BUCKET = os.getenv("OSS_BUCKET", "bucket")

# --- 初始化 S3 客户端 ---
s3_client = boto3.client(
    's3',
    endpoint_url=OSS_URL,
    aws_access_key_id=OSS_ACCESS_KEY,
    aws_secret_access_key=OSS_SECRET_KEY,
    config=Config(signature_version='s3v4')
)

# --- 缓存清理逻辑 ---
# 指定一个统一的临时文件夹前缀，方便清理
TEMP_PREFIX = "ppt_parse_"

def clean_up_old_temps():
    """
    启动时清理所有以 TEMP_PREFIX 开头的残留临时文件夹
    """
    root_temp = tempfile.gettempdir()
    pattern = os.path.join(root_temp, f"{TEMP_PREFIX}*")
    folders = glob.glob(pattern)
    
    for folder in folders:
        try:
            if os.path.isdir(folder):
                shutil.rmtree(folder)
                print(f"Cleanup: Removed old temp folder {folder}")
            else:
                os.remove(folder)
        except Exception as e:
            print(f"Cleanup Error: Failed to remove {folder}, error: {e}")

@app.on_event("startup")
async def startup_event():
    # 程序启动时执行清理逻辑
    print("Service starting... Cleaning up old temporary files.")
    clean_up_old_temps()

class ProcessRequest(BaseModel):
    object_name: str

# --- 工具函数 ---
def get_safe_ext_and_mime(file_obj, original_filename):
    """
    结合 Magic Number 和后缀名判断最准确的 MIME 和后缀
    """
    file_obj.seek(0)
    head = file_obj.read(2048)
    file_obj.seek(0)

    mime_type = magic.from_buffer(head, mime=True)
    orig_ext = os.path.splitext(original_filename)[1].lstrip('.').lower()
    
    if mime_type.startswith("text/") or mime_type == "application/octet-stream":
        final_ext = orig_ext if orig_ext else "bin"
        final_mime, _ = mimetypes.guess_type(f"file.{final_ext}")
        final_mime = final_mime or mime_type
    else:
        detected_ext = mimetypes.guess_extension(mime_type)
        if detected_ext:
            final_ext = detected_ext.lstrip('.')
        else:
            final_ext = orig_ext if orig_ext else "bin"
        final_mime = mime_type

    return final_ext, final_mime

def calculate_stream_md5(file_obj) -> str:
    """分块计算 MD5，不占用大内存"""
    hasher = hashlib.md5()
    file_obj.seek(0)
    while True:
        chunk = file_obj.read(1024 * 1024) 
        if not chunk:
            break
        hasher.update(chunk)
    file_obj.seek(0) 
    return hasher.hexdigest()

# --- 主接口 ---
@app.post("/ppt/url")
async def process_ppt(request: ProcessRequest):
    object_name = request.object_name
    result_data = {"filenames": [], "text": []}
    
    file_ext = os.path.splitext(object_name.lower())[1]
    if file_ext not in ['.pptx', '.ppt']:
        return {"filenames": [object_name], "text": []}

    # 使用 prefix=TEMP_PREFIX 确保启动清理逻辑能识别该目录
    with tempfile.TemporaryDirectory(prefix=TEMP_PREFIX) as tmp_dir:
        download_path = os.path.join(tmp_dir, f"source{file_ext}")
        
        try:
            # 1. 下载原始文件
            response = s3_client.get_object(Bucket=OSS_BUCKET, Key=object_name)
            with open(download_path, "wb") as f:
                for chunk in response['Body'].iter_chunks(chunk_size=1024 * 1024):
                    f.write(chunk)
            
            # 2. 如果是 .ppt，转换为 .pptx (LibreOffice)
            target_pptx = download_path
            if file_ext == '.ppt':
                try:
                    subprocess.run(
                        ["soffice", "--headless", "--convert-to", "pptx", download_path, "--outdir", tmp_dir],
                        check=True,
                        capture_output=True,
                        timeout=60
                    )
                    # 转换后的文件名固定为 source.pptx
                    target_pptx = os.path.join(tmp_dir, "source.pptx")
                except Exception as e:
                    return {"filenames": [object_name], "text": []}

            # 3. 解析文本内容
            try:
                prs = Presentation(target_pptx)
                for slide in prs.slides:
                    slide_text = [
                        shape.text.strip() 
                        for shape in slide.shapes 
                        if hasattr(shape, "text") and shape.text.strip()
                    ]
                    result_data["text"].append("\n".join(slide_text))
            except Exception as e:
                result_data["text"].append(f"Text Parse Error: {str(e)}")

            # 4. 提取媒体文件并上传
            if zipfile.is_zipfile(target_pptx):
                with zipfile.ZipFile(target_pptx, 'r') as zip_ref:
                    for file_info in zip_ref.infolist():
                        if file_info.filename.startswith('ppt/media/') and not file_info.filename.endswith('/'):
                            # 在 tmp_dir 中创建临时媒体文件，方便后续清理
                            with tempfile.NamedTemporaryFile(dir=tmp_dir, delete=True) as tmp_media:
                                with zip_ref.open(file_info.filename) as source:
                                    shutil.copyfileobj(source, tmp_media)
                                tmp_media.flush()
                                
                                extension, content_type = get_safe_ext_and_mime(tmp_media, file_info.filename)
                                md5_hash = calculate_stream_md5(tmp_media)
                                new_key = f"{md5_hash}.{extension}"
                                
                                if new_key not in result_data["filenames"]:
                                    file_exists = False
                                    try:
                                        s3_client.head_object(Bucket=OSS_BUCKET, Key=new_key)
                                        file_exists = True
                                    except:
                                        file_exists = False

                                    if not file_exists:
                                        tmp_media.seek(0)
                                        s3_client.put_object(
                                            Bucket=OSS_BUCKET, Key=new_key,
                                            Body=tmp_media, ContentType=content_type
                                        )
                                    result_data["filenames"].append(new_key)

            return result_data

        except Exception as e:
            raise HTTPException(status_code=500, detail=f"Internal Error processing {object_name}: {str(e)}")


@app.get("/health", tags=["Health"])
async def health_check():
    """
    基础健康检查，用于 K8S 的 Liveness Probe
    """
    return {"status": "healthy", "service": "ppt-parse-service"}